import urllib as a
import sys
#url = "https://s3-us-west-2.amazonaws.com/tarunproject/acne.txt";
#url = raw_input("enter url");
url = "https://s3-us-west-2.amazonaws.com/tarunproject/sample.txt";
b = a.urlopen(url); #getting file
c = b.read(); #reading the file and string into a variable
#c = "An Object is an entity that is represented by information in a collection or database. User queries are coordinated against the database data.Be that as it may, instead of traditional SQL questions of a database, in Information Retrieval the outcomes returned might possibly may not match query, so results are normally ranked. This ranking of results is a key distinction of Information Retrieval searching contrasted with database searching.";
c = c.replace("\r","").replace("\r","").replace("\n",""); #replacing tab and new line chars
import re 
#d = re.findall("[a-zA-Z0-9!? ,'""*;:-]*[^.\\r\\n,]",c);
d = c.split("."); #splitting the document by .
e = [];
for i in d:
	e.append(i.strip());
#print e;
#print f;
from sklearn.feature_extraction.text import TfidfVectorizer #importing dependecies
from sklearn.metrics.pairwise import linear_kernel
tfidf = TfidfVectorizer(); #creating a TfidfVectorizer object
try:
	tf = tfidf.fit_transform(e)  
	mat = linear_kernel(tf,tf) #creating a matrix of frequency vectors
except:
	print("sorry");
	sys.exit();	
import numpy
m = numpy.array(mat).tolist(); #converting that to a list
#print m;
sum = [];
for i in m: #looping throught the values and adding all the rows to get the resultant
    sum1 = 0;
    for j in i:
        sum1 = sum1+j;
    sum.append(sum1); #appending sum to sum matrix  
#print sum;
val = [];
#print numpy.median(sum),numpy.max(sum);
#print len(sum),len(m);
for i in range(len(sum)): #if we want sentences greter than median use this 
    if(sum[i] >= round(numpy.median(sum))):
        val.append(i);
#print val;
dict = {};
for i in sum:#creating a dictionary with keys and sums and values as documents
	if(dict.has_key(i) == False):
		dict[i]=[];
for i in range(len(sum)):
	dict[sum[i]].append(i);
userinput  = input("Enter number of sentences max is "+str(len(sum)));#taking use input
finallist = []; 
dectoace = sorted(dict.keys());
shit = dectoace;
shit.reverse();
#print shit;
for keys in shit: #appending values to the final list
	if(len(finallist) <=  userinput):
		for i in dict[keys]:
			finallist.append(i);
di = [];
index =1;
for i in finallist: #creating the sentences
	di.append(str(index)+": "+d[i].strip());
	index+=1;
from pptx import Presentation #impoting ppt conversion library
prs = Presentation() #creating presentation object

for i in di[0:userinput]: #creating the ppt of all senteces with one setence in each slide
	print i;
	title_slide_layout = prs.slide_layouts[0];
	slide = prs.slides.add_slide(title_slide_layout);
	title = slide.shapes.title;
	subtitle = slide.placeholders[1];
	subtitle.text="";
	title.text = i;
prs.save('test.pptx')


	








